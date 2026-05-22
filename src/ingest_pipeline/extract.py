import re
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

import anthropic
import fitz  # pymupdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from .config import CHUNK_SIZE_BYTES, HAIKU_MODEL, MAX_TOKENS
from .api import call_claude, call_claude_vision
from .state import ManifestFile, State, append_journal, journal_event

_FR_WORDS: frozenset[str] = frozenset({
    "les", "des", "est", "une", "qui", "que", "dans", "sur", "pour", "avec",
    "pas", "au", "aux", "par", "ce", "je", "tu", "il", "nous", "vous", "ils",
    "mais", "ou", "et", "donc", "or", "ni", "car", "du", "de", "la", "le",
    "en", "un", "plus", "bien", "tout", "même", "très", "aussi", "comme",
    "être", "avoir", "faire", "cette", "leur", "leurs", "dont", "dont",
})

_EN_WORDS: frozenset[str] = frozenset({
    "the", "and", "for", "that", "this", "with", "from", "are", "have", "not",
    "you", "all", "can", "its", "your", "our", "they", "was", "will", "been",
    "has", "had", "would", "about", "which", "their", "when", "there", "also",
    "into", "more", "some", "than", "then", "time", "what", "said", "each",
    "other", "just", "know", "take", "only", "come", "most", "over", "such",
})


def detect_locale(text: str) -> str:
    """Detect dominant language (fr/en/und) from text using stopword frequency."""
    words = {w.lower() for w in text.split() if w.isalpha() and len(w) > 1}
    fr_score = len(words & _FR_WORDS)
    en_score = len(words & _EN_WORDS)
    if fr_score == 0 and en_score == 0:
        return "und"
    return "fr" if fr_score > en_score else "en"


def _rel(path: str) -> str:
    """Return path relative to CWD; fall back to original if not relative."""
    try:
        return str(Path(path).relative_to(Path.cwd()))
    except ValueError:
        return path


def _chunk_text(text: str, chunk_size: int) -> list[str]:
    """Split text into chunks of at most chunk_size bytes on line boundaries."""
    lines = text.splitlines()
    chunks: list[str] = []
    current: list[str] = []
    current_size = 0
    for line in lines:
        line_bytes = len((line + "\n").encode("utf-8"))
        if current_size + line_bytes > chunk_size and current:
            chunks.append("\n".join(current))
            current, current_size = [line], line_bytes
        else:
            current.append(line)
            current_size += line_bytes
    if current:
        chunks.append("\n".join(current))
    return chunks if chunks else [""]


def has_complexity_signals(text: str) -> bool:
    lines = text.splitlines()

    # Pipe tables
    if any("|" in line and line.strip().startswith("|") for line in lines):
        return True

    # Aligned columns: two or more consecutive spaces appearing mid-line
    aligned_count = sum(1 for line in lines if re.search(r"\S {3,}\S", line))
    if aligned_count >= 3:
        return True

    # Repeated non-trivial lines (≥3 occurrences of same non-empty line)
    non_trivial = [line.strip() for line in lines if len(line.strip()) > 10]
    counts = Counter(non_trivial)
    if any(count >= 3 for count in counts.values()):
        return True

    return False


def _build_front_matter(record: ManifestFile, source_path: Path, locale: str = "und") -> str:
    ingest_id = record.id.removeprefix("sha256:")
    extracted_at = datetime.now(timezone.utc).isoformat()
    return "\n".join([
        "---",
        "source: " + _rel(record.source_path),
        "category: " + record.category,
        "locale: " + locale,
        "extracted_at: " + extracted_at,
        "ingest_id: " + ingest_id,
        "---",
        "",
        "",
    ])


def extract_text_file(
    source_path: Path,
    dest_path: Path,
    record: ManifestFile,
    prompts: dict[str, str],
    client: anthropic.Anthropic,
    rpm: int,
    state: State,
    dest_root: Path,
) -> None:
    try:
        raw_text = source_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            raw_text = source_path.read_text(encoding="latin-1")
        except UnicodeDecodeError:
            raw_text = source_path.read_text(encoding="utf-8", errors="replace")

    if len(raw_text.encode("utf-8")) > 500_000:
        chunks = _chunk_text(raw_text, 500_000)
        processed_parts = []
        for chunk in chunks:
            user_content = "Source: " + record.source_path + "\n\nContent:\n" + chunk
            processed_parts.append(call_claude(
                client,
                HAIKU_MODEL,
                prompts["extract_text"],
                user_content,
                MAX_TOKENS["extract"],
                rpm,
                state,
                dest_root,
            ))

        content = "\n\n".join(processed_parts)
    elif has_complexity_signals(raw_text):
        user_content = "Source: " + record.source_path + "\n\nContent:\n" + raw_text
        content = call_claude(
            client,
            HAIKU_MODEL,
            prompts["extract_text"],
            user_content,
            MAX_TOKENS["extract"],
            rpm,
            state,
            dest_root,
        )
    else:
        content = raw_text

    locale = detect_locale(raw_text)
    front_matter = _build_front_matter(record, source_path, locale)
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    dest_path.write_text(front_matter + content, encoding="utf-8")


def extract_pdf(source_path: Path) -> str:
    """Extract text from PDF using pymupdf with header/footer stripping."""
    doc = fitz.open(source_path)

    # Collect all lines per page to detect headers/footers
    page_lines: list[list[str]] = []
    for page in doc:
        text = page.get_text("text")
        lines = [line.strip() for line in text.splitlines()]
        page_lines.append(lines)

    # Identify repeated lines (headers/footers) appearing in >=80% of pages
    # Only applies when there are 2+ pages (need repetition to detect headers/footers)
    total_pages = len(page_lines)
    lines_to_strip: set[str] = set()

    if total_pages >= 2:
        threshold = total_pages * 0.8  # Use float comparison, not int

        # Count non-empty line occurrences across all pages
        line_counts: dict[str, int] = {}
        for lines in page_lines:
            seen_in_page = set()
            for line in lines:
                if line and line not in seen_in_page:
                    seen_in_page.add(line)
                    line_counts[line] = line_counts.get(line, 0) + 1

        # Identify lines to strip (appear on >=80% of pages)
        lines_to_strip = {line for line, count in line_counts.items() if count >= threshold}

    # Extract text with stripping
    all_text = []
    for lines in page_lines:
        filtered_lines = [line for line in lines if line and line not in lines_to_strip]
        all_text.append("\n".join(filtered_lines))

    doc.close()
    return "\n\n".join(all_text)


def extract_docx(source_path: Path) -> str:
    """Extract text from DOCX using python-docx, excluding headers/footers."""
    doc = Document(source_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n\n".join(paragraphs)


def extract_pptx(source_path: Path, warnings: list[dict]) -> str:
    """Extract text from PPTX using python-pptx, including speaker notes.

    Image-only slides are recorded as dicts in warnings (caller writes journal events).
    """
    prs = Presentation(source_path)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Extract text from shapes
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text.strip()
                if text:
                    slide_text_parts.append(text)

        # Extract speaker notes
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # If slide has no text at all, record warning for caller
        if not slide_text_parts and not notes_text:
            warnings.append({"reason": "image-only slide", "slide": slide_num})
            continue

        # Build slide content
        slide_content = "\n\n".join(slide_text_parts)
        if notes_text:
            slide_content += "\n\n## Speaker Notes\n" + notes_text

        slides_text.append(slide_content)

    return "\n\n".join(slides_text)


def extract_xlsx(source_path: Path) -> str:
    """Extract text from XLSX using openpyxl as markdown tables per sheet."""
    wb = load_workbook(source_path, data_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        # Get all rows with data
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # Build markdown table
        table_lines = ["## " + sheet_name, ""]

        # Convert rows to strings, handling None values
        for row in rows:
            cells = [str(cell) if cell is not None else "" for cell in row]
            table_lines.append("| " + " | ".join(cells) + " |")

        sheets_text.append("\n".join(table_lines))

    wb.close()
    return "\n\n".join(sheets_text)


def extract_binary_doc(
    source_path: Path,
    dest_path: Path,
    record: ManifestFile,
    prompts: dict[str, str],
    client: anthropic.Anthropic,
    rpm: int,
    state: State,
    dest_root: Path,
) -> None:
    """Extract text from binary documents (PDF, DOCX, PPTX, XLSX) and always call Claude."""
    suffix = source_path.suffix.lower()

    # Extract raw text based on format
    if suffix == ".pdf":
        raw_text = extract_pdf(source_path)
    elif suffix == ".docx":
        raw_text = extract_docx(source_path)
    elif suffix == ".pptx":
        pptx_warnings: list[dict] = []
        raw_text = extract_pptx(source_path, pptx_warnings)
        for w in pptx_warnings:
            append_journal(
                dest_root,
                journal_event("file_skipped", step=1, source=str(source_path), cmd=state.command, **w),
            )
    elif suffix == ".xlsx":
        raw_text = extract_xlsx(source_path)
    else:
        raise ValueError("Unsupported binary document format: " + suffix)

    # For PDF, handle chunking if text > CHUNK_SIZE_BYTES
    if suffix == ".pdf" and len(raw_text.encode("utf-8")) > CHUNK_SIZE_BYTES:
        chunks = _chunk_text(raw_text, CHUNK_SIZE_BYTES)
        processed_parts = []
        for chunk in chunks:
            user_content = "Source: " + record.source_path + "\n\nContent:\n" + chunk
            content = call_claude(
                client,
                HAIKU_MODEL,
                prompts["extract_text"],
                user_content,
                MAX_TOKENS["extract"],
                rpm,
                state,
                dest_root,
            )
            processed_parts.append(content)

        final_content = "\n\n".join(processed_parts)
    else:
        # Always call Claude for binary docs
        user_content = "Source: " + record.source_path + "\n\nContent:\n" + raw_text
        final_content = call_claude(
            client,
            HAIKU_MODEL,
            prompts["extract_text"],
            user_content,
            MAX_TOKENS["extract"],
            rpm,
            state,
            dest_root,
        )

    locale = detect_locale(raw_text)
    front_matter = _build_front_matter(record, source_path, locale)
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    dest_path.write_text(front_matter + final_content, encoding="utf-8")


_IMAGE_MEDIA_TYPES: dict[str, str] = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


def extract_image(
    source_path: Path,
    dest_path: Path,
    record: ManifestFile,
    prompts: dict[str, str],
    client: anthropic.Anthropic,
    rpm: int,
    state: State,
    dest_root: Path,
) -> None:
    suffix = source_path.suffix.lower()
    media_type = _IMAGE_MEDIA_TYPES[suffix]
    image_data = source_path.read_bytes()
    front_matter = _build_front_matter(record, source_path)
    content = call_claude_vision(
        client,
        HAIKU_MODEL,
        prompts["extract_image"],
        image_data,
        media_type,
        "Source: " + record.source_path,
        MAX_TOKENS["image"],
        rpm,
        state,
        dest_root,
    )
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    dest_path.write_text(front_matter + content, encoding="utf-8")


def extract_file(
    record: ManifestFile,
    prompts: dict[str, str],
    client: anthropic.Anthropic,
    rpm: int,
    state: State,
    dest_root: Path,
) -> None:
    source_path = Path(record.source_path)
    dest_path = Path(record.destination_path)

    if record.category == "text":
        extract_text_file(source_path, dest_path, record, prompts, client, rpm, state, dest_root)
    elif record.category == "binary-doc":
        extract_binary_doc(source_path, dest_path, record, prompts, client, rpm, state, dest_root)
    elif record.category == "binary-image":
        extract_image(source_path, dest_path, record, prompts, client, rpm, state, dest_root)
    else:
        raise ValueError("Unknown category: " + record.category)


def run_extract_step(
    records: list[ManifestFile],
    prompts: dict[str, str],
    client: anthropic.Anthropic,
    rpm: int,
    state: State,
    dest_root: Path,
    on_file: callable = None,
) -> list[str]:
    errors: list[str] = []
    for record in records:
        dest_path = Path(record.destination_path)
        if dest_path.exists():
            continue
        if (
            record.category in ("binary-doc", "binary-image")
            and state.max_binary_mb > 0
            and record.size_bytes > state.max_binary_mb * 1024 * 1024
        ):
            record.status = "skipped-oversized"
            append_journal(
                dest_root,
                journal_event(
                    "file_skipped",
                    step=1,
                    source=record.source_path,
                    reason="oversized",
                    size_bytes=record.size_bytes,
                    cmd=state.command,
                ),
            )
            continue
        if on_file:
            on_file(record.source_path)
        try:
            extract_file(record, prompts, client, rpm, state, dest_root)
            record.status = "extracted"
            state.completed_files.append(record.source_path)
        except NotImplementedError as e:
            record.status = "skipped"
            errors.append(str(record.source_path) + ": " + str(e))
        except Exception as e:
            record.status = "failed"
            state.failed_files.append(record.source_path)
            errors.append(str(record.source_path) + ": " + str(e))
            append_journal(dest_root, journal_event("file_error", file=record.source_path, error=str(e), cmd=state.command))
    return errors
